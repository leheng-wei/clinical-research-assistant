import streamlit as st
import fitz
import hashlib
import requests
import os
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from docx import Document
from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
from io import BytesIO
from datetime import datetime, timedelta
from tenacity import retry, stop_after_attempt, wait_exponential
import pandas as pd
import base64
import json
import uuid
import io
import csv
import re

# ===== 配置 =====
# 
# ===== 配置 =====
DEEPSEEK_API_KEY = st.secrets.get("DEEPSEEK_API_KEY", "")
if not DEEPSEEK_API_KEY:
    st.error("请在 .streamlit/secrets.toml 中设置 DEEPSEEK_API_KEY")

# 常量定义
MAX_FILE_SIZE = 200 * 1024 * 1024  # 200MB
MAX_FILES_PER_BATCH = 5
UPLOAD_DIR = "uploaded_pdfs"
LOGO_PATH = "bofu_logo.png"  # 移动到配置部分
os.makedirs(UPLOAD_DIR, exist_ok=True)

# ===== 文件处理函数 =====
def validate_file(uploaded_file):
    """验证上传文件的大小和类型"""
    if len(uploaded_file.getvalue()) > MAX_FILE_SIZE:
        st.error(f"文件大小超过限制（200MB）：{uploaded_file.name}")
        return False
    if not uploaded_file.name.lower().endswith('.pdf'):
        st.error(f"只支持PDF文件：{uploaded_file.name}")
        return False
    return True

def sanitize_filename(filename):
    """清理文件名，移除不安全的字符"""
    safe_chars = set("abcdefghijklmnopqrstuvwxyzABCDEFGHIJKLMNOPQRSTUVWXYZ0123456789-._")
    filename = ''.join(c for c in filename if c in safe_chars)
    return filename or "unnamed_file"

def process_large_file(uploaded_file):
    """分块处理大文件以优化内存使用"""
    chunk_size = 1024 * 1024  # 1MB chunks
    file_bytes = uploaded_file.getvalue()
    total_chunks = (len(file_bytes) + chunk_size - 1) // chunk_size
    
    progress_bar = st.progress(0)
    text_blocks = []
    
    for i in range(0, len(file_bytes), chunk_size):
        chunk = file_bytes[i:i + chunk_size]
        doc = fitz.open(stream=chunk, filetype="pdf")
        
        for page in doc:
            blocks = page.get_text("blocks")
            for b in blocks:
                if b[6] == 0:
                    clean_text = b[4].replace('\x00', '').strip()
                    if clean_text and len(clean_text) > 20:
                        text_blocks.append(clean_text)
        
        progress = (i + chunk_size) / len(file_bytes)
        progress_bar.progress(min(progress, 1.0))
    
    return "\n".join(text_blocks)

def process_file_with_status(uploaded_file):
    """添加处理状态提示的文件处理函数"""
    status_placeholder = st.empty()
    progress_placeholder = st.empty()
    progress_bar = progress_placeholder.progress(0)
    
    try:
        # 步骤1：验证文件
        status_placeholder.info("📑 验证文件...")
        if not validate_file(uploaded_file):
            status_placeholder.error("❌ 文件验证失败")
            return None
        progress_bar.progress(25)
        
        # 步骤2：提取文本
        status_placeholder.info("📖 提取文本...")
        file_bytes = uploaded_file.getvalue()
        file_hash = hashlib.md5(file_bytes).hexdigest()
        full_text = extract_text_from_pdf_by_hash(file_hash, file_bytes)
        if not full_text:
            status_placeholder.error("❌ 文本提取失败")
            return None
        progress_bar.progress(50)
        
        # 步骤3：分析内容
        status_placeholder.info("🤖 分析内容...")
        trimmed_text = full_text[:30000]
        text_hash = hashlib.md5(trimmed_text.encode("utf-8")).hexdigest()
        result = extract_study_by_hash(text_hash, trimmed_text)
        if not result:
            status_placeholder.error("❌ 内容分析失败")
            return None
        progress_bar.progress(75)
        
        # 步骤4：生成报告
        status_placeholder.info("📊 生成报告...")
        try:
            # 生成CSV
            csv_lines = []
            for line in result.splitlines():
                if "|" in line and not line.startswith("|---") and not line.lower().startswith("| 要素"):
                    parts = [p.strip() for p in line.strip("|").split("|")]
                    if len(parts) >= 2:
                        csv_lines.append(",".join(parts))
            csv_lines.insert(0, "要素,内容")
            
            # 生成PPT
            prs = Presentation()
            prs.slide_width = Inches(13.33)
            prs.slide_height = Inches(7.5)
            
            # 生成Word
            word_bytes = generate_word_table(csv_lines, "博扶AI创意组", uploaded_file.name)
            
            progress_bar.progress(100)
            status_placeholder.success("✅ 处理完成")
            return result, csv_lines, prs, word_bytes
            
        except Exception as e:
            status_placeholder.error("❌ 报告生成失败")
            return None
            
    except Exception as e:
        status_placeholder.error("❌ 处理失败")
        return None
    finally:
        # 清理状态显示
        progress_placeholder.empty()
        status_placeholder.empty()

# ===== 文本提取函数 =====
@st.cache_data(show_spinner="📖 正在解析 PDF", max_entries=100)
def extract_text_from_pdf_by_hash(file_hash: str, file_bytes: bytes):
    """根据文件大小选择处理方式"""
    if len(file_bytes) > MAX_FILE_SIZE:
        return process_large_file(file_bytes)
    else:
        doc = fitz.open(stream=file_bytes, filetype="pdf")
        text_blocks = []
        for page in doc:
            blocks = page.get_text("blocks")
            for b in blocks:
                if b[6] == 0:
                    clean_text = b[4].replace('\x00', '').strip()
                    if clean_text and len(clean_text) > 20:
                        text_blocks.append(clean_text)
        return "\n".join(text_blocks)

@st.cache_data(show_spinner="🤖 正在调用模型...", max_entries=100)
def extract_study_by_hash(text_hash: str, text: str):
    return extract_study_design(text)

@retry(stop=stop_after_attempt(3), wait=wait_exponential(multiplier=1, min=4, max=10))
def extract_study_design(text):
    """提取研究设计信息，带重试机制"""
    prompt = f"""作为临床研究专家，请从以下文献中提取结构化研究设计信息，并输出为 Markdown 表格，包含以下字段：

| 要素 | 内容 |
|------|------|
| 研究类型 | RCT / 队列 / 病例对照 / 横断面等 |
| 是否多中心 | 是/否 |
| 是否盲法 | 单盲 / 双盲 / 开放标签 |
| 纳入/排除标准 | 简要列出 |
| 干预措施（实验组） | 药物、剂量、频率等 |
| 干预措施（对照组） | 如安慰剂/标准治疗等 |
| 患者人数 | 样本总量及分组数量 |
| 主要终点指标 | 疗效指标名称与评估方式 |
| 次要/其他终点指标 | 包括所有非主要终点的疗效观察指标，如 secondary outcomes、exploratory outcomes、其他效果评估等 |
| 关键量化指标 | 所有被用于量化分析、机制研究、建模或分组的重要变量。不限终点、不限疗效相关，如 insulin resistance, clearance 等 |
| 安全性终点指标 | 不良事件、实验室指标等 |
| 统计分析方法 | 所用统计工具和模型 |
| 临床试验注册号 | 如有请列出 |

文献内容如下：
{text}
"""
    try:
        response = requests.post(
            "https://api.deepseek.com/v1/chat/completions",
            headers={"Authorization": f"Bearer {DEEPSEEK_API_KEY}"},
            json={
                "model": "deepseek-chat",
                "messages": [{"role": "user", "content": prompt}],
                "temperature": 0,
                "max_tokens": 4000
            },
            timeout=90
        )
        response.raise_for_status()
        return response.json()["choices"][0]["message"]["content"]
    except requests.exceptions.RequestException as e:
        st.error(f"模型调用失败（将自动重试）：{e}")
        raise
    except Exception as e:
        st.error(f"处理过程中发生错误：{e}")
        return ""

def extract_supplementary_notes(result):
    """提取补充说明内容"""
    if not isinstance(result, str):
        return None
    
    # 使用更精确的方式检测补充说明
    supplementary_pattern = r"补充说明[：:]\s*(.*?)(?=\n\n|\Z)"
    match = re.search(supplementary_pattern, result, re.DOTALL)
    
    if match:
        return match.group(1).strip()
    return None

# ===== 文档生成函数 =====
def generate_word_table(csv_lines, team_name, source_file):
    doc = Document()
    heading = doc.add_heading(f"{team_name} · 临床研究结构化提取报告", 0)
    heading.alignment = WD_PARAGRAPH_ALIGNMENT.CENTER
    subtitle = doc.add_paragraph(f"📄 来源文献：{source_file}", style="Intense Quote")
    subtitle.alignment = WD_PARAGRAPH_ALIGNMENT.CENTER
    doc.add_paragraph()
    table = doc.add_table(rows=0, cols=2)
    table.style = 'Table Grid'
    for row in csv_lines[1:]:
        if "," in row:
            key, value = row.split(",", 1)
            cells = table.add_row().cells
            cells[0].text = key.strip()
            cells[1].text = value.strip()
    for row in table.rows:
        for cell in row.cells:
            for para in cell.paragraphs:
                para.paragraph_format.line_spacing = 1.5
                para.paragraph_format.space_after = Pt(6)
                para.paragraph_format.alignment = WD_PARAGRAPH_ALIGNMENT.LEFT
                para.runs[0].font.size = Pt(11)
    doc.add_paragraph(f"导出时间：{datetime.now().strftime('%Y-%m-%d %H:%M:%S')}")
    word_stream = BytesIO()
    doc.save(word_stream)
    word_stream.seek(0)
    return word_stream

# ===== 历史记录管理函数 =====
def load_history():
    """从本地文件加载历史记录"""
    history_file = 'history.json'
    if not os.path.exists(history_file):
        print("历史记录文件不存在，将创建新的记录文件")
        return []
    
    try:
        with open(history_file, 'r', encoding='utf-8') as f:
            data = json.load(f)
            if not isinstance(data, dict) or 'records' not in data:
                print("历史记录文件格式错误，将重置记录")
                return []
            records = data.get('records', [])
            if not isinstance(records, list):
                print("历史记录格式错误，将重置记录")
                return []
            
            # 清理超过1周的记录
            current_time = datetime.now()
            one_week_ago = current_time - timedelta(days=7)
            cleaned_records = [
                record for record in records 
                if datetime.strptime(record['时间'], '%Y-%m-%d %H:%M:%S') > one_week_ago
            ]
            return cleaned_records
    except Exception as e:
        print(f"加载历史记录时发生错误: {str(e)}")
        return []

def save_history(records):
    """保存历史记录到本地文件"""
    if not isinstance(records, list):
        print("记录格式错误，取消保存")
        return False
    
    try:
        # 清理超过1周的记录
        current_time = datetime.now()
        one_week_ago = current_time - timedelta(days=7)
        cleaned_records = [
            record for record in records 
            if datetime.strptime(record['时间'], '%Y-%m-%d %H:%M:%S') > one_week_ago
        ]
        
        data = {
            "records": cleaned_records,
            "metadata": {
                "last_updated": datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
                "total_records": len(cleaned_records),
                "cleaned_records": len(records) - len(cleaned_records)
            }
        }
        with open('history.json', 'w', encoding='utf-8') as f:
            json.dump(data, f, ensure_ascii=False, indent=2)
        return True
    except Exception as e:
        print(f"保存历史记录时发生错误: {str(e)}")
        return False

def delete_record(record_id):
    """删除单条历史记录"""
    if not record_id:
        print("记录ID无效")
        return False
        
    if 'history' not in st.session_state:
        print("历史记录状态未初始化")
        return False
        
    try:
        records = st.session_state.history
        original_length = len(records)
        st.session_state.history = [r for r in records if r.get('id') != record_id]
        
        if len(st.session_state.history) == original_length:
            print(f"未找到ID为 {record_id} 的记录")
            return False
            
        if save_history(st.session_state.history):
            print(f"成功删除记录 {record_id}")
            return True
        else:
            print("删除记录后保存失败")
            return False
    except Exception as e:
        print(f"删除记录时发生错误: {str(e)}")
        return False

# ===== 页面展示函数 =====
def show_history():
    st.subheader("📜 处理历史记录")
    if len(st.session_state.history) > 0:
        history_df = pd.DataFrame(st.session_state.history)
        history_df['时间'] = pd.to_datetime(history_df['时间'])
        history_df.sort_values('时间', ascending=False, inplace=True)

        # 展示历史记录表格
        st.dataframe(history_df)

        # 提供下载历史记录功能
        csv = history_df.to_csv(index=False).encode('utf-8')
        st.download_button(
            label="下载历史记录 CSV",
            data=csv,
            file_name="历史记录.csv",
            mime="text/csv"
        )

        # 提供清除历史记录按钮
        if st.button("清除所有历史记录"):
            st.session_state.history.clear()
            st.success("历史记录已清除！")
    else:
        st.warning("没有历史记录。")

def clear_uploaded_files():
    """清理上传的文件"""
    if "uploaded_pdfs" in os.listdir():
        files = os.listdir("uploaded_pdfs")
        if len(files) > 0:
            for file in files:
                os.remove(os.path.join("uploaded_pdfs", file))
            st.success("上传文件已清理！")
        else:
            st.warning("没有上传文件可清理！")
    else:
        st.warning("文件目录不存在！")

def admin_dashboard():
    st.title("后台管理 - 临床研究设计结构化助手")
    
    # 获取历史记录，读取本地文件或者缓存数据
    if 'history' not in st.session_state:
        st.session_state.history = []
    
    action = st.selectbox("选择操作", ["查看历史记录", "清理上传文件", "退出后台"])

    if action == "查看历史记录":
        show_history()
    elif action == "清理上传文件":
        clear_uploaded_files()
    else:
        st.write("退出后台管理界面")

# ===== 主程序 =====
st.set_page_config(layout="wide", page_title="临床研究设计结构化助手")
page = st.radio("选择页面", ["主页", "后台管理"], horizontal=True)

# 初始化或加载历史记录
if 'history' not in st.session_state:
    st.session_state.history = load_history()

if page == "主页":
    # 上传文献并展示结构化提取结果
    st.markdown('<div style="display: flex; align-items: flex-start; padding-top: 12px; font-size: 45px;">🤖</div>', unsafe_allow_html=True)
    st.markdown('<h1 style="margin: 0;">临床研究设计结构化助手</h1>', unsafe_allow_html=True)
    st.caption("💡 自动识别 PDF 文献中的研究设计信息，支持导出为 CSV / PPT / Word")
    
    # 添加处理状态说明
    with st.expander("💡 处理说明", expanded=False):
        st.info("""
        **处理步骤：**
        1. 📑 验证文件
        2. 📖 提取文本
        3. 🤖 分析内容
        4. 📊 生成报告
        
        **提示：**
        - 每篇文献处理时间约1-2分钟
        - 处理过程中请勿关闭页面
        """)
    
    uploaded_files = st.file_uploader("📄 上传PDF文件（支持多选，每个文件限制200MB）", type=["pdf"], accept_multiple_files=True)
    total_files = len(uploaded_files) if uploaded_files else 0
    if uploaded_files:
       total_files = len(uploaded_files)
    if total_files > 5:
        st.error(f"❌ 超出单次处理限制（5篇）")
        # 只处理前5篇
        current_batch = uploaded_files[:5]
        queued_files = uploaded_files[5:]
    else:
        current_batch = uploaded_files
        queued_files = []
    if current_batch:
        tabs = st.tabs([f"📄 {i+1}. {file.name}" for i, file in enumerate(current_batch)])
        for idx, (tab, uploaded_file) in enumerate(zip(tabs, current_batch)):
            with tab:
            # 处理逻辑
              ...
    else:
        st.info("请上传一篇或多篇 PDF 文献以开始处理。")

    st.caption(f"📚 当前处理：{len(current_batch)} 篇文献" + (f" | 队列中：{len(queued_files)} 篇" if queued_files else ""))
    
    # 处理上传的文献
    tabs = st.tabs([f"📄 {i+1}. {file.name}" for i, file in enumerate(current_batch)])

    for idx, (tab, uploaded_file) in enumerate(zip(tabs, current_batch)):
        with tab:
            # 计算文件的哈希值
            file_bytes = uploaded_file.read()
            uploaded_file.seek(0)  # 将文件指针重置到文件开头，以便后续处理
            file_hash = hashlib.sha256(file_bytes).hexdigest()

            # 使用新的处理函数
            result = process_file_with_status(uploaded_file)
            
            if not result:
                st.error(f"❌ 处理失败：{uploaded_file.name}")
                continue

            result, csv_lines, prs, word_bytes = result
            
            # 结构化转 CSV 行
            today = datetime.now().strftime("%Y%m%d")
            csv_bytes = BytesIO("\n".join(csv_lines).encode("utf-8"))
            csv_bytes.seek(0)

            # ===== PPT 生成（带 LOGO 封面） =====
            prs = Presentation()
            prs.slide_width = Inches(13.33)
            prs.slide_height = Inches(7.5)
            
            # 封面页
            cover = prs.slides.add_slide(prs.slide_layouts[6])
            title_box = cover.shapes.add_textbox(Inches(1), Inches(0.8), Inches(11), Inches(1.5))
            tf = title_box.text_frame
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = "研究设计提取报告"
            run.font.size = Pt(42)
            run.font.name = "微软雅黑"
            run.font.bold = True
            run.font.color.rgb = RGBColor(0, 51, 102)
            p.alignment = PP_ALIGN.CENTER
            
            if os.path.exists(LOGO_PATH):
                cover.shapes.add_picture(LOGO_PATH, Inches(4.8), Inches(2.0), height=Inches(0.6))
            
            sub_box = cover.shapes.add_textbox(Inches(1), Inches(4.0), Inches(11), Inches(0.8))
            tf2 = sub_box.text_frame
            tf2.text = f"源文件：  {uploaded_file.name}"
            tf2.paragraphs[0].font.size = Pt(20)
            tf2.paragraphs[0].font.name = "微软雅黑"
            tf2.paragraphs[0].alignment = PP_ALIGN.CENTER
            
            # 内容页
            for i, row in enumerate(csv_lines[1:]):
                parts = row.split(",", 1)
                if len(parts) == 2:
                    slide = prs.slides.add_slide(prs.slide_layouts[1])
                    title_shape = slide.shapes.title
                    title_shape.text = parts[0]
                    title_frame = title_shape.text_frame
                    title_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
                    title_run = title_frame.paragraphs[0].runs[0]
                    title_run.font.name = "微软雅黑"
                    title_run.font.bold = True
                    title_run.font.color.rgb = RGBColor(0, 51, 102)
                    
                    textbox = slide.placeholders[1]
                    textbox.text = parts[1].replace("；", "\n")
                    for p in textbox.text_frame.paragraphs:
                        for run in p.runs:
                            run.font.size = Pt(18)
                            run.font.name = "微软雅黑"
                    
                    footer = slide.shapes.add_textbox(Inches(0.5), Inches(6.9), Inches(12), Inches(0.5))
                    tf_footer = footer.text_frame
                    tf_footer.text = f"博扶AI创意组 · 结构化助手 · 第 {i+1} 页"
                    tf_footer.paragraphs[0].font.size = Pt(10)
                    tf_footer.paragraphs[0].font.name = "微软雅黑"
                    tf_footer.paragraphs[0].alignment = PP_ALIGN.RIGHT
            
            # 添加补充说明页
            if "补充说明" in result:
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                title_shape = slide.shapes.title
                title_shape.text = "补充说明"
                title_frame = title_shape.text_frame
                title_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
                title_run = title_frame.paragraphs[0].runs[0]
                title_run.font.name = "微软雅黑"
                title_run.font.bold = True
                title_run.font.color.rgb = RGBColor(0, 51, 102)
                
                textbox = slide.placeholders[1]
                textbox.text = result.split("补充说明：")[1].strip()
                for p in textbox.text_frame.paragraphs:
                    for run in p.runs:
                        run.font.size = Pt(18)
                        run.font.name = "微软雅黑"
                
                footer = slide.shapes.add_textbox(Inches(0.5), Inches(6.9), Inches(12), Inches(0.5))
                tf_footer = footer.text_frame
                tf_footer.text = f"博扶AI创意组 · 结构化助手 · 补充说明页"
                tf_footer.paragraphs[0].font.size = Pt(10)
                tf_footer.paragraphs[0].font.name = "微软雅黑"
                tf_footer.paragraphs[0].alignment = PP_ALIGN.RIGHT
            
            pptx_bytes = BytesIO()
            prs.save(pptx_bytes)
            pptx_bytes.seek(0)

            # ===== 展示结果 =====
            st.success("✅ 已成功提取结构化研究设计信息")
            st.markdown(result.strip(), unsafe_allow_html=True)
            
            # 添加提示信息
            st.info("💡 提示：请及时下载生成的文件，历史记录将在7天后自动清理。")
            st.warning("⚠️ 注意：文件仅保存在浏览器会话中，关闭页面后将无法访问。")

            # 生成文件名
            today = datetime.now().strftime("%Y%m%d")
            base_name = os.path.splitext(uploaded_file.name)[0]

            st.markdown("#### 📁 下载导出文件")
            col1, col2, col3 = st.columns(3)
            with col1:
                st.download_button("📥 下载 CSV", csv_bytes, f"{today}_{base_name}_结构化.csv", mime="text/csv")
            with col2:
                st.download_button("📊 下载 PPT", pptx_bytes, f"{today}_{base_name}_结构化.pptx", mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")
            with col3:
                st.download_button("📄 下载 Word", word_bytes, f"{today}_{base_name}_结构化.docx", mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document")

            # 保存记录到历史
            record = {
                "id": str(uuid.uuid4()),
                "文件名": uploaded_file.name,
                "hash": file_hash,
                "时间": datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
                "提取内容": result.strip(),
                "文件数据": {
                    "CSV": base64.b64encode(csv_bytes.getvalue()).decode('utf-8'),
                    "PPT": base64.b64encode(pptx_bytes.getvalue()).decode('utf-8'),
                    "Word": base64.b64encode(word_bytes.getvalue()).decode('utf-8')
                }
            }
            st.session_state.history.append(record)
            save_history(st.session_state.history)

            st.markdown("---")
            st.subheader("📜 历史处理记录")
            for record in st.session_state.history:
                with st.expander(f"📄 `{record['文件名']}`"):
                    col1, col2 = st.columns([0.9, 0.1])
                    with col1:
                        st.markdown(f"📁 文件 Hash: `{record['hash']}`")
                        st.markdown(f"⏰ 时间: {record['时间']}")
                        st.markdown(f"📄 提取内容:\n{record['提取内容']}")

                        # 从base64字符串恢复数据并提供下载
                        st.markdown(f"#### 下载文件:")
                        download_col1, download_col2, download_col3 = st.columns(3)
                        with download_col1:
                            try:
                                csv_data = base64.b64decode(record['文件数据']['CSV'])
                                st.download_button("📥 下载 CSV", 
                                                data=csv_data, 
                                                file_name=f"{record['文件名']}_结构化.csv",
                                                mime="text/csv",
                                                key=f"csv_{record['id']}")
                            except Exception as e:
                                st.error(f"CSV数据加载失败")
                        
                        with download_col2:
                            try:
                                pptx_data = base64.b64decode(record['文件数据']['PPT'])
                                st.download_button("📊 下载 PPT", 
                                                data=pptx_data,
                                                file_name=f"{record['文件名']}_结构化.pptx",
                                                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                                                key=f"ppt_{record['id']}")
                            except Exception as e:
                                st.error(f"PPT数据加载失败")
                        
                        with download_col3:
                            try:
                                word_data = base64.b64decode(record['文件数据']['Word'])
                                st.download_button("📄 下载 Word",
                                                data=word_data,
                                                file_name=f"{record['文件名']}_结构化.docx",
                                                mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                                                key=f"word_{record['id']}")
                            except Exception as e:
                                st.error(f"Word数据加载失败")
                        
                    with col2:
                        if st.button("🗑️", key=f"delete_{record['id']}", help="删除此记录"):
                            delete_record(record['id'])
                            st.rerun()

elif page == "后台管理":
    admin_dashboard()

st.markdown("---")
st.markdown("© 2025 博扶AI创意组 · 医学文献结构化助手")
